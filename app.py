import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
from gtts import gTTS
import tempfile
import os

# --- INSTÄLLNINGAR ---
st.set_page_config(page_title="Jag Lär Mig", page_icon="📖", layout="wide")

# --- SESSIONS-HANTERING ---
if "subjects" not in st.session_state:
    st.session_state.subjects = {"Allmänt": ""}
if "current_subject" not in st.session_state:
    st.session_state.current_subject = "Allmänt"

# --- FUNKTIONER ---

# (Övriga funktioner för extrahera text, gTTS, etc., är oförändrade)

def extract_text_from_pdf(pdf_file):
    text = ""
    reader = PdfReader(pdf_file)
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def generate_speech_simple(text):
    try:
        tts = gTTS(text=text, lang='sv')
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            tts.save(fp.name)
            return fp.name
    except Exception as e:
        st.error(f"Kunde inte skapa ljud: {e}")
        return None

def get_gemini_response(prompt, context, api_key):
    if not api_key:
        return "⚠️ Fel: API-nyckel saknas. Lägg in nyckeln i Secrets!"
    
    genai.configure(api_key=api_key) 
    
    system_instruction = (
        "Du är en smart och pedagogisk studiecoach i appen 'Jag Lär Mig'. "
        "Din uppgift är att hjälpa användaren att förstå sitt studiematerial. "
        "Var tydlig, uppmuntrande och svara alltid på svenska."
    )
    model = genai.GenerativeModel('gemini-1.5-pro', system_instruction=system_instruction)
    
    full_prompt = f"Studiematerial:\n{context}\n\nUppgift/Fråga: {prompt}"
    
    try:
        return model.generate_content(full_prompt).text
    except Exception as e:
        error_msg = str(e)
        
        # Mer detaljerad felhantering
        if "API key not valid" in error_msg:
            st.error("❌ Google avvisar nyckeln! Kontrollera att den är korrekt i Secrets.")
        elif "NotFound" in error_msg:
            st.error("❌ Modellen hittades inte. Kontrollera att du använder rätt modellnamn.")
        elif "quota" in error_msg.lower():
            st.error("⚠️ Du har nått din kvot hos Google AI. Vänta eller uppgradera din plan.")
        elif "timeout" in error_msg.lower():
            st.error("⏳ Anropet tog för lång tid. Testa igen senare.")
        else:
            st.error(f"🚨 Oväntat fel: {error_msg}")
        
        return f"Ett fel uppstod vid AI-anropet.\n\nDetaljer: {error_msg}"



# --- SIDOPANEL (MENY) ---

# --- NY HANTERING AV API-NYCKELN ---
with st.sidebar:
    st.title("📖 Jag Lär Mig")
    
    # Försök hämta nyckeln från Secrets
    if "GEMINI_API_KEY" in st.secrets:
        # Om nyckeln finns i secrets, använd den.
        api_key = st.secrets["GEMINI_API_KEY"]
        st.success("🔑 Nyckel laddad från Secrets!")
    else:
        # Annars, visa instruktioner
        st.warning("⚠️ NYCKEL SAKNAS. Se nedan hur du lägger in den.")
        api_key = "" # Håll nyckeln tom
        
        st.info("""
        **För att fixa felet permanent:**
        1. Gå till Streamlit Cloud dashboard.
        2. Klicka på **Secrets**.
        3. Lägg till nyckeln under namnet **GEMINI_API_KEY**.
        4. Starta om appen.
        """)

    st.divider()

    # --- RESTEN AV SIDAN (som tidigare) ---
    st.subheader("📂 Mina Ämnen")
    # ... resten av koden är densamma ...
    subject_list = list(st.session_state.subjects.keys())
    
    # Välj ämne
    selected_sub = st.selectbox("Välj ämne att plugga:", subject_list, index=subject_list.index(st.session_state.current_subject))
    st.session_state.current_subject = selected_sub
    
    # Skapa nytt ämne
    new_sub = st.text_input("Lägg till nytt ämne (t.ex. Kemi):")
    if st.button("Skapa mapp") and new_sub:
        st.session_state.subjects[new_sub] = ""
        st.session_state.current_subject = new_sub
        st.success(f"Mappen '{new_sub}' skapad!")
        st.rerun()

    st.divider()
    
    # Uppladdning
    st.subheader(f"📥 Ladda upp till: {st.session_state.current_subject}")
    uploaded_files = st.file_uploader("Släpp filer här (PDF, PPTX)", accept_multiple_files=True)
    
    if st.button("Spara materialet"):
        text_data = st.session_state.subjects[st.session_state.current_subject]
        count = 0
        for file in uploaded_files:
            if file.name.endswith(".pdf"):
                text_data += f"\n--- {file.name} ---\n" + extract_text_from_pdf(file)
                count += 1
            elif file.name.endswith(".pptx"):
                text_data += f"\n--- {file.name} ---\n" + extract_text_from_pptx(file)
                count += 1
        
        st.session_state.subjects[st.session_state.current_subject] = text_data
        st.success(f"Sparade {count} filer i {st.session_state.current_subject}!")

# --- HUVUDVY ---
st.header(f"Studerar: {st.session_state.current_subject}")

current_material = st.session_state.subjects[st.session_state.current_subject]

if not current_material:
    st.info("👈 Den här mappen är tom. Börja med att ladda upp material i menyn!")
else:
    # --- FLIKAR ---
    tab1, tab2, tab3 = st.tabs(["📝 Material & Struktur", "🎧 Lyssna", "💬 Förhör & Chatt"])

    # FLIK 1: REDIGERA
    with tab1:
        st.subheader("Ditt material")
        st.caption("Här kan du se texten som appen läst in och ändra om något blev fel.")
        
        # Redigeringsfönster
        edited_text = st.text_area("Innehåll", current_material, height=300)
        
        if st.button("Spara ändringar i texten"):
            st.session_state.subjects[st.session_state.current_subject] = edited_text
            st.success("Uppdaterat!")
            st.rerun()

        st.divider()
        if st.button("✨ Dela upp texten i kapitel (AI)"):
            with st.spinner("Analyserar struktur..."):
                chapters = get_gemini_response(
                    "Dela upp texten i tydliga kapitel med rubriker.", 
                    edited_text, api_key
                )
                st.markdown(chapters)

    # FLIK 2: LYSSNA
    with tab2:
        st.subheader("Uppläsning")
        
        text_to_read = st.text_area("Text att läsa upp:", value=edited_text[:3000], height=150)

        if st.button("▶️ Spela upp"):
            with st.spinner("Skapar ljud..."):
                audio_path = generate_speech_simple(text_to_read)
                if audio_path:
                    st.audio(audio_path, format="audio/mp3")

    # FLIK 3: CHATT / FÖRHÖR
    with tab3:
        st.subheader("Plugga med AI")
        
        # Förslag på knappar
        c1, c2, c3 = st.columns(3)
        if c1.button("Skapa ett prov"):
            with st.spinner("Skapar prov..."):
                test = get_gemini_response("Skapa ett prov med 5 frågor + facit.", edited_text, api_key)
                st.markdown(test)
        
        if c2.button("Sammanfatta allt"):
            with st.spinner("Sammanfattar..."):
                summary = get_gemini_response("Sammanfatta det viktigaste i punktform.", edited_text, api_key)
                st.markdown(summary)

        # Chatt
        st.divider()
        user_q = st.chat_input("Ställ en fråga om materialet...")
        if user_q:
            st.chat_message("user").write(user_q)
            with st.spinner("Tänker..."):
                ans = get_gemini_response(user_q, edited_text, api_key)
                st.chat_message("assistant").write(ans)
